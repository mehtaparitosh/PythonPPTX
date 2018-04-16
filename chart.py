# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.chart.data import BubbleChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
from pptx.chart.axis import MajorGridlines
from pptx.util import Inches, Pt, Cm, Px, Mm
from pptx.chart.data import ChartData
from pptx.chart.data import XyChartData
from pptx.chart.data import BubbleChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.enum.chart import XL_TICK_MARK
from pptx.dml.color import RGBColor
from pptx.dml.color import ColorFormat
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.oxml.xmlchemy import OxmlElement
import enum
import mysql.connector
from mysql.connector import errorcode
import pandas as pd
from decimal import Decimal
import pptx
import pptx.util
import glob
import scipy.misc
from pptx.oxml.xmlchemy import OxmlElement
from pptx.enum.text import PP_ALIGN



# create presentation with 1 slide ------
prs = Presentation()
array = []

def populate_data():
	print array[0][0]
	print array[1][0]


def iter_row(cursor,size=2):
	while True:
		rows = cursor.fetchmany(size)
		if not rows:
			break
		for row in rows:

			yield row
	pass


def Query_data(cnx):

	try:
		cursor = cnx.cursor()
		cursor.execute("SELECT * FROM entries")
		# row = cursor.fetchone()

		for row in iter_row(cursor,2):
			array.append(row)
			print(row)


	except Error as e:
		print(e)
	finally:
		cursor.close()

def connect():

	try:
	  cnx = mysql.connector.connect(user='root',password='password@123',host='127.0.0.1',
	                                database='invescoapp_development')
	except mysql.connector.Error as err:
	  if err.errno == errorcode.ER_ACCESS_DENIED_ERROR:
	    print("Something is wrong with your user name or password")
	  elif err.errno == errorcode.ER_BAD_DB_ERROR:
	    print("Database does not exist")
	  else:
	    print(err)
	else:
		print "Success!"
		if cnx.is_connected():
			print('Connected to  MYSQL database')
			Query_data(cnx)

	finally :
		print "Closing connection."
		cnx.close()

def iter_cells(table):
    for row in table.rows:
        for cell in row.cells:
            yield cell

def format_table_headers(table):
	_set_cell_border(table.cell(0,0))
	table.cell(0,0).fill.solid()
	table.cell(0,0).fill.fore_color.rgb = 	RGBColor(0xff, 0xff, 0xff)
	# table.first_row.fill.fore_color.rgb = RGBColor(0xff, 0xff, 0xff)
	cols = 6
	for c in range(cols):
		_set_cell_border(table.cell(0,c))
		table.cell(0,c).fill.solid()
		table.cell(0,c).fill.fore_color.rgb = RGBColor(0xff, 0xff, 0xff)



	for pp in table.cell(0,0).text_frame.paragraphs:
		for run in pp.runs:
			run.font.bold = True

	# table.cell(0,0).text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
# cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
	table.cell(0,0).text_frame.vertical_anchor = 4
# cell.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.BOTTOM

def rule_based_formatting(table):
	class ColorCode(enum.Enum):
		high = 1
		medium = 2
		low = 3
	row = 14
	for r in range(row):
		table.cell(r,0).fill.solid()
		table.cell(r,0).fill.fore_color.rgb = 	RGBColor(154,205,50)

	table.cell(0,0).fill.solid()
	table.cell(0,0).fill.fore_color.rgb = RGBColor(0xff, 0xff, 0xff)

	pass

def rule_based_formatting_2(table):
	class ColorCode(enum.Enum):
		high = 1
		medium = 2
		low = 3
	row = 14
	for r in range(len(table.rows)):
		table.cell(r,0).fill.solid()
		table.cell(r,0).fill.fore_color.rgb = 	RGBColor(255,191,0)

	table.cell(0,0).fill.solid()
	table.cell(0,0).fill.fore_color.rgb = RGBColor(0xff, 0xff, 0xff)
	pass

def rule_based_formatting_3(table):
	class ColorCode(enum.Enum):
		high = 1
		medium = 2
		low = 3
	row = 14
	for r in range(len(table.rows)):
		table.cell(r,0).fill.solid()
		table.cell(r,0).fill.fore_color.rgb = 	RGBColor(254,0,5)

	table.cell(0,0).fill.solid()
	table.cell(0,0).fill.fore_color.rgb = RGBColor(0xff, 0xff, 0xff)
	pass

def rule_based_formatting_4(table):
	class ColorCode(enum.Enum):
		high = 1
		medium = 2
		low = 3
	row = 14
	for r in range(len(table.rows)):
		table.cell(r,0).fill.solid()
		table.cell(r,0).fill.fore_color.rgb = 	RGBColor(154,205,50)

	table.cell(0,0).fill.solid()
	table.cell(0,0).fill.fore_color.rgb = RGBColor(0xff, 0xff, 0xff)
	pass

def color_red(table,row,column):
	table.cell(row,column).fill.solid()
	table.cell(row,column).fill.fore_color.rgb = 	RGBColor(254,0,5)

def color_yellow(table,row,column):
	table.cell(row,column).fill.solid()
	table.cell(row,column).fill.fore_color.rgb = 	RGBColor(255,191,0)
	pass
def color_green(table,row,column):
	table.cell(row,column).fill.solid()
	table.cell(row,column).fill.fore_color.rgb = 	RGBColor(154,205,50)
	pass

def rule_based_formatting_inventory(table):
	col = 6
	c= 1
	while (c < col):
		if(array[0][c] == None):
			pass
		elif(array[0][c]<80):
			color_red(table,c+1,3)
		elif(array[0][c]>90):
			color_green(table,c+1,3)
		else:
			color_yellow(table,c+1,3)

		c = c+1
	c = 1
	while (c < col):
		if(array[1][c]==None):
			pass
		elif(array[1][c]<80):
			color_red(table,c+1,4)
		elif(array[1][c]>90):
			color_green(table,c+1,4)
		else:
			color_yellow(table,c+1,4)

		c = c+1
	pass

def rule_based_formatting_security_baseline(table):
	col = 11
	c = 7
	while (c < col ):
		if(array[0][c] == None):
			pass
		elif (array[0][c] <60):
			color_red(table,c+2,3)
		elif (array[0][c]>80 ):
			color_green(table,c+2,3)
		else:
			color_yellow(table,c+2,3)
		c = c+1

	c = 7
	while (c < col ):
		if(array[1][c] == None):
			pass
		elif (array[1][c] <60):
			color_red(table,c+2,4)
		elif (array[1][c]>80 ):
			color_green(table,c+2,4)
		else:
			color_yellow(table,c+2,4)
		c = c+1
	pass

def rule_based_formatting_security_baseline_2(table):
	col = 15
	c = 11
	while (c < col ):
		if(array[0][c] == None):
			pass
		elif (array[0][c] <60):
			color_red(table,c-10,3)
		elif (array[0][c]>80 ):
			color_green(table,c-10,3)
		else:
			color_yellow(table,c-10,3)
		c = c+1
		pass

	c = 11
	while (c < col ):
		if(array[1][c] == None ):
			pass
		elif (array[1][c] <60):
			color_red(table,c-10,4)
		elif (array[1][c]>80 ):
			color_green(table,c-10,4)
		else:
			color_yellow(table,c-10,4)
		c = c+1
		pass
	pass

def rule_based_formatting_vulnerabilitylandscape(table):
	col = 23
	c = 18
	while(c < col ):
		if(array[0][c] == None):
			pass
		elif(array[0][c] > 20 ):
			color_red(table,c-9,3)
		elif (array[0][c] < 10 ):
			color_green(table,c-9,3)
		else:
			color_yellow(table,c-9,3)
		c = c + 1
	c = 18
	while(c < col ):
		if(array[1][c] == None):
			pass
		elif(array[1][c] > 20 ):
			color_red(table,c-9,4)
		elif (array[1][c] < 10 ):
			color_green(table,c-9,4)
		else:
			color_yellow(table,c-9,4)
		c = c + 1
	pass

def rule_based_formatting_vulnerabilitylandscape2(table):
	col = 28
	c = 23
	while(c < col ):
		if(array[0][c] == None):
			pass
		elif(array[0][c] > 90 ):
			color_red(table,c-21,3)
		elif(array[0][c]< 60 ):
			color_green(table,c-21,3)
		else :
			color_yellow(table,c-21,3)
		c = c+1

	c = 23
	while(c < col ):
		if(array[1][c] == None):
			pass
		elif(array[1][c] > 90 ):
			color_red(table,c-21,4)
		elif(array[1][c]< 60 ):
			color_green(table,c-21,4)
		else :
			color_yellow(table,c-21,4)
		c = c+1
	pass

def rule_based_formatting_vulnerabilitylandscape3(table):
	col = 33
	c = 28
	while ( c < col ):
		if(array[0][c] == None):
			pass
		elif (array[0][c] > 180 ):
			color_red(table,c-20,3)
		elif(array[0][c] < 90 ):
			color_green(table,c-20,3)
		else:
			color_yellow(table,c-20,3)
		c = c+1
	c = 28
	while ( c < col ):
		if(array[1][c] == None):
			pass
		elif (array[1][c] > 180 ):
			color_red(table,c-20,4)
		elif(array[1][c] < 90 ):
			color_green(table,c-20,4)
		else:
			color_yellow(table,c-20,4)
		c = c+1
	pass

def rule_based_formatting_monitoring(table):
	col = 38
	c = 34
	while(c < col ):
		if(array[0][c] == None):
			pass
		elif(array[0][c] > 15) :
			color_red(table,c-33,3)
		elif(array[0][c] < 5 ):
			color_green(table,c-33,3)
		else:
			color_yellow(table,c-33,3)
		c = c + 1

	c = 34
	while(c < col ):
		if(array[1][c] == None):
			pass
		elif(array[1][c] > 15) :
			color_red(table,c-33,4)
		elif(array[1][c] < 5 ):
			color_green(table,c-33,4)
		else:
			color_yellow(table,c-33,4)
		c = c + 1
	pass

def color_in_grey(table,r,c_start,c_end):
	diff = c_end - c_start + 1
	for col in range(diff):
		table.cell(r,col+c_start).fill.solid()
		table.cell(r,col+c_start).fill.fore_color.rgb = RGBColor(220,220,220)
	pass

def format_table_specific_content(table):
	color_in_grey(table,1,3,5)
	color_in_grey(table,8,3,5)
	color_in_grey(table,13,3,5)
	pass
def format_table_specific_content_2(table):
	color_in_grey(table,8,3,5)
	# color_in_grey(table,8,3,5)
	# color_in_grey(table,13,3,5)
	pass

def format_table_specific_content_3(table):
	color_in_grey(table,1,3,5)
	color_in_grey(table,7,3,5)
	color_in_grey(table,13,3,5)
	pass

def format_table_specific_content_4(table):

	pass

def format_table_content(table):
	row = 14
	r = 1
	for r in range(len(table.rows)):
		table.cell(r,1).fill.solid()
		table.cell(r,1).fill.fore_color.rgb = RGBColor(211, 211, 211)


	for r in range(len(table.rows)):
		table.cell(r,2).fill.solid()
		table.cell(r,2).fill.fore_color.rgb = RGBColor(220, 220, 220)

	# color_in_grey(table,1,3,5)
	# color_in_grey(table,8,3,5)
	# color_in_grey(table,13,3,5)




	table.cell(0,1).fill.solid()
	table.cell(0,1).fill.fore_color.rgb = RGBColor(0xff, 0xff, 0xff)
	table.cell(0,2).fill.solid()
	table.cell(0,2).fill.fore_color.rgb = RGBColor(0xff, 0xff, 0xff)



	pass
def format_table_content_es(table):

	pass

def fixed_table_content_es(table):
	table.cell(1,0).text = "INV"
	table.cell(1,1).text = "Asset Inventory"
	table.cell(6,0).text = "SEB"
	table.cell(6,1).text = "Security Baseline"
	table.cell(10,0).text = "VLL"
	table.cell(10,1).text = "Vulnerability Landscape"
	table.cell(15,0).text = "MON"
	table.cell(15,1).text = "Security Monitoring"
	table.cell(16,0).text = "IAM"
	table.cell(16,1).text = "Privileged Access"






	table.cell(1,3).text = " "
	table.cell(1,4).text = " "
	table.cell(1,5).text = " "



	table.cell(1,3).text = "Networks"
	table.cell(2,3).text = "Servers "
	table.cell(3,3).text = "Endpoints"
	table.cell(4,3).text = "Database"
	table.cell(5,3).text = "Applications"
	table.cell(1,4).text = "Jason L"
	table.cell(2,4).text = "Harris T"
	table.cell(3,4).text = "Ashley T"
	table.cell(4,4).text = "Peter C"
	table.cell(5,4).text = "Sue L"
	table.cell(6,3).text = "Networks"
	table.cell(7,3).text = "Servers "
	table.cell(8,3).text = "Endpoints"
	table.cell(9,3).text = "Database"
	table.cell(6,4).text = "Cory J"
	table.cell(7,4).text = "Harris T"
	table.cell(8,4).text = "Wes M"
	table.cell(9,4).text = "Peter C"
	table.cell(10,3).text = "Networks"
	table.cell(11,3).text = "Servers "
	table.cell(12,3).text = "Endpoints"
	table.cell(13,3).text = "Database"
	table.cell(14,3).text = "Applications"
	table.cell(10,4).text = "Pavan D"
	table.cell(11,4).text = "Yogesh S"
	table.cell(12,4).text = "Lloyd E"
	table.cell(13,4).text = "Srinivasulu P"
	table.cell(14,4).text = "Sue L"
	table.cell(15,3).text = "Overall"
	table.cell(16,3).text = "Overall"
	table.cell(15,4).text = "Jos V"
	table.cell(16,4).text = "Brian S"


	for cell in iter_cells(table):
		for paragraph in cell.text_frame.paragraphs:
			for run in paragraph.runs:
				run.font.size = Pt(9)
				run.font.color.rgb = RGBColor(0,0,0)
				run.font.bold = True
	r = 1
	while(r < 10):
		color_green(table,r,0)
		r = r + 1

	r = 10
	while( r < 17):
		color_red(table,r,0)
		r = r + 1
	r = 1
	while(r < 17):
		color_in_grey(table,r,1,2)
		r = r + 1





	pass

def fixed_table_content(table):
	table.cell(1,0).text = "Inventory"
	table.cell(1,1).text = "Accuracy and completeness of asset information"
	table.cell(1,2).text = " % of inventory completion (Owner / IT Custodian must be filled to be considered a complete inventory item.)"
	table.cell(1,3).text = " "
	table.cell(1,4).text = " "
	table.cell(1,5).text = " "

	table.cell(2,2).text = "Networks"

	count = 1 
	while(count < 7 ):
		if(array[0][count]!= None):
			table.cell(count+1,3).text = str(Decimal(array[0][count]).normalize())+ "%"
		else : 
			table.cell(count+1,3).text = " "
		count = count + 1 
	count = 1 
	while(count < 7 ):
		if(array[1][count]!= None):
			table.cell(count+1,4).text = str(Decimal(array[1][count]).normalize())+ "%"
		else : 
			table.cell(count+1,4).text = " "
		count = count + 1 

	

	table.cell(3,2).text = "Servers"


	table.cell(4,2).text = "Endpoints"
	table.cell(5,2).text = "Database"
	table.cell(6,2).text = "Applications"
	rule_based_formatting_inventory(table)
	table.cell(7,2).text = "# of unknown Assets (Those which are discovered but not found in the “source of truth” probably CMDB)"
	table.cell(8,0).text = "Security Baseline"
	table.cell(8,1).text = "Alignment with industry best practices"
	table.cell(8,2).text = "% of assets with security baseline (# of assets as a factor of whole estate, for which the baseline is Agreed. (Coverage)"
	table.cell(8,3).text = " "
	table.cell(8,4).text = " "

	table.cell(9,2).text = "Networks"

	count = 7 
	while(count < 11 ):
		if(array[0][count]!= None):
			table.cell(count+2,3).text = str(Decimal(array[0][count]).normalize())+ "%"
		else : 
			table.cell(count+2,3).text = " "
		count = count + 1 
	count = 7
	while(count < 11 ):
		if(array[1][count]!= None):
			table.cell(count+2,4).text = str(Decimal(array[1][count]).normalize())+ "%"
		else : 
			table.cell(count+2,4).text = " "
		count = count + 1 

	table.cell(10,2).text = "Servers"
	table.cell(11,2).text = "Endpoints"
	table.cell(12,2).text = "Database"
	table.cell(13,2).text = " "
	rule_based_formatting_security_baseline(table)






	run = table.cell(1, 0).text_frame.paragraphs[0].runs[0]
	run.font.bold = True

	run = table.cell(1, 1).text_frame.paragraphs[0].runs[0]
	run.font.bold = True

	run = table.cell(8, 0).text_frame.paragraphs[0].runs[0]
	run.font.bold = True

	run = table.cell(8, 1).text_frame.paragraphs[0].runs[0]
	run.font.bold = True
	pass

def fixed_table_content_2(table):
	table.cell(1,0).text = "Security Baseline"
	table.cell(1,1).text = "Alignment with industry best practices"
	table.cell(1,3).text = " "
	table.cell(1,4).text = " "
	table.cell(1,5).text = " "
	table.cell(5,0).text = "Identity and Access Management"
	table.cell(5,1).text = "Measure of Privileged accounts"
	table.cell(8,0).text = "Vulnerability Landscape"
	table.cell(8,1).text = "Measure of Vulnerabilities in infrastructure "
	table.cell(8,3).text = " "
	table.cell(8,4).text = " "
	table.cell(8,5).text = " "
	table.cell(9,2).text = "Networks"


	count = 18
	while(count < 23 ):
		if(array[0][count]!= None):
			table.cell(count-9,3).text = str(Decimal(array[0][count]).normalize())+ "%"
		else : 
			table.cell(count-9,3).text = " "
		count = count + 1 
	count = 18
	while(count < 23 ):
		if(array[1][count]!= None):
			table.cell(count-9,4).text = str(Decimal(array[1][count]).normalize())+ "%"
		else : 
			table.cell(count-9,4).text = " "
		count = count + 1 

	table.cell(10,2).text = "Servers"
	table.cell(11,2).text = "Endpoints"
	table.cell(12,2).text = "Database"
	table.cell(13,2).text = "Applications"
	rule_based_formatting_vulnerabilitylandscape(table)


	count = 11 
	while(count < 15 ):
		if(array[0][count]!= None):
			table.cell(count-10,3).text = str(Decimal(array[0][count]).normalize())+ "%"
		else : 
			table.cell(count-10,3).text = " "
		count = count + 1 
	count = 11 
	while(count < 15 ):
		if(array[1][count]!= None):
			table.cell(count-10,4).text = str(Decimal(array[1][count]).normalize())+ "%"
		else : 
			table.cell(count-10,4).text = " "
		count = count + 1 

	table.cell(1,2).text = "Networks"
	table.cell(2,2).text = "Servers"
	table.cell(3,2).text = "Endpoints"
	table.cell(4,2).text = "Database"
	rule_based_formatting_security_baseline_2(table)
	table.cell(5,2).text = "Total # of Admin accounts"
	table.cell(5,3).text = str(array[0][15])
	table.cell(5,4).text = str(array[1][15])

	table.cell(6,2).text = "# of 2FA Admin accounts "
	table.cell(6,3).text = str(array[0][16])
	table.cell(6,4).text = str(array[1][16])
	table.cell(7,2).text = "# of Generic Admin accounts without owner"
	table.cell(7,3).text = str(array[0][17])
	table.cell(7,4).text = str(array[1][17])
	table.cell(8,2).text = "% of assets with High and critical (4&5) vulnerabilities"





	run = table.cell(1, 0).text_frame.paragraphs[0].runs[0]
	run.font.bold = True

	run = table.cell(1, 1).text_frame.paragraphs[0].runs[0]
	run.font.bold = True

	run = table.cell(5, 0).text_frame.paragraphs[0].runs[0]
	run.font.bold = True

	run = table.cell(5, 1).text_frame.paragraphs[0].runs[0]
	run.font.bold = True

	run = table.cell(8, 0).text_frame.paragraphs[0].runs[0]
	run.font.bold = True

	run = table.cell(8, 1).text_frame.paragraphs[0].runs[0]
	run.font.bold = True




	pass

def fixed_table_content_3(table):
	table.cell(1,0).text = "Vulnerability Landscape"
	table.cell(1,1).text = "Are we in Control ?"
	table.cell(1,3).text = " "
	table.cell(1,4).text = " "
	table.cell(1,5).text = " "

	table.cell(7,1).text = "What is our Risk ?"
	table.cell(7,2).text = "Average Age of Open Vulnerabilities (in # of Days)"
	table.cell(13,0).text = "Monitoring"
	table.cell(13,1).text = "Measure of assets not monitored"
	table.cell(13,2).text = "% of Assets not being monitored"


	# table.cell(8,0).text = "Vulnerability Landscape"
	# table.cell(8,1).text = "Measure of Vulnerabilities in infrastructure "
	table.cell(7,3).text = " "
	table.cell(7,4).text = " "
	table.cell(7,5).text = " "





	table.cell(1,2).text = "Average Time to Remediate last cycle (in # of Days)"
	table.cell(2,2).text = "Networks"
	count = 23
	while(count < 28 ):
		if(array[0][count]!= None):
			table.cell(count-21,3).text = str(Decimal(array[0][count]).normalize())+ "Days"
		else : 
			table.cell(count-21,3).text = " "
		count = count + 1 
	count = 23 
	while(count < 28 ):
		if(array[1][count]!= None):
			table.cell(count-21,4).text = str(Decimal(array[1][count]).normalize())+ "Days"
		else : 
			table.cell(count-21,4).text = " "
		count = count + 1 
	table.cell(3,2).text = "Servers"
	table.cell(4,2).text = "Endpoints"
	table.cell(5,2).text = "Database"
	table.cell(6,2).text = "Applications"
	rule_based_formatting_vulnerabilitylandscape2(table)
	table.cell(8,2).text = "Networks"
	count = 28
	while(count < 34 ):
		if(array[0][count]!= None):
			table.cell(count-20,3).text = str(Decimal(array[0][count]).normalize())+ "Days"
		else : 
			table.cell(count-20,3).text = " "
		count = count + 1 
	count = 28
	while(count < 34 ):
		if(array[1][count]!= None):
			table.cell(count-20,4).text = str(Decimal(array[1][count]).normalize())+ "Days"
		else : 
			table.cell(count-20,4).text = " "
		count = count + 1 

	table.cell(9,2).text = "Servers"
	table.cell(10,2).text = "Endpoints"
	table.cell(11,2).text = "Database"
	table.cell(12,2).text = "Applications"

	rule_based_formatting_vulnerabilitylandscape3(table)
	table.cell(14,2).text = "Networks"

	val = array[0][34]
	val2 = array[1][34]
	if (val == None ):
		pass
	elif ( val > 15):
		color_red(table,14,3)
	elif(val < 5 ):
		color_green(table,14,3)
	else :
		color_yellow(table,14,3)

	if(val2 == None ):
		pass
	elif ( val2 > 15):
		color_red(table,14,4)
	elif(val2 < 5 ):
		color_green(table,14,4)
	else :
		color_yellow(table,14,4)











	run = table.cell(1, 0).text_frame.paragraphs[0].runs[0]
	run.font.bold = True

	run = table.cell(1, 1).text_frame.paragraphs[0].runs[0]
	run.font.bold = True

	# run = table.cell(5, 0).text_frame.paragraphs[0].runs[0]
	# run.font.bold = True

	run = table.cell(7, 1).text_frame.paragraphs[0].runs[0]
	run.font.bold = True

	run = table.cell(13, 0).text_frame.paragraphs[0].runs[0]
	run.font.bold = True

	run = table.cell(13, 1).text_frame.paragraphs[0].runs[0]
	run.font.bold = True




	pass

	pass

def fixed_table_content_4(table):
	table.cell(1,0).text = "Monitoring"
	table.cell(1,1).text = "Measure of assets not monitored"
	table.cell(1,2).text = "Servers"
	count = 34 
	while(count < 38 ):
		if(array[0][count]!= None):
			table.cell(count-33,3).text = str(Decimal(array[0][count]).normalize())+ "%"
		else : 
			table.cell(count-33,3).text = " "
		count = count + 1 
	count = 34 
	while(count < 38 ):
		if(array[1][count]!= None):
			table.cell(count-33,4).text = str(Decimal(array[1][count]).normalize())+ "%"
		else : 
			table.cell(count-33,4).text = " "
		count = count + 1 
	# table.cell(1,3).text = str(Decimal(array[0][34]).normalize())+"%"
	# table.cell(1,4).text = str(Decimal(array[1][34]).normalize())+"%"
	table.cell(2,2).text = "Endpoints"
	# table.cell(2,3).text = str(Decimal(array[0][35]).normalize())+"%"
	# table.cell/(2,4).text = str(Decimal(array[1][35]).normalize())+"%"
	table.cell(3,2).text = "Database"
	# table.cell(3,3).text = str(Decimal(array[0][36]).normalize())+"%"
	# table.cell(3,4).text = str(Decimal(array[1][36]).normalize())+"%"
	table.cell(4,2).text = "Applications"
	# table.cell(4,3).text = str(Decimal(array[0][37]).normalize())+"%"
	# table.cell(4,4).text = str(Decimal(array[1][37]).normalize())+"%"
	rule_based_formatting_monitoring(table)

	run = table.cell(1, 0).text_frame.paragraphs[0].runs[0]
	run.font.bold = True

	run = table.cell(1, 1).text_frame.paragraphs[0].runs[0]
	run.font.bold = True

	pass

def SubElement(parent, tagname, **kwargs):
        element = OxmlElement(tagname)
        element.attrib.update(kwargs)
        parent.append(element)
        return element

def _set_cell_border(cell, border_color="000000", border_width='12700'):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for lines in ['a:lnL','a:lnR','a:lnT','a:lnB']:
        ln = SubElement(tcPr, lines, w=border_width, cap='flat', cmpd='sng', algn='ctr')
        solidFill = SubElement(ln, 'a:solidFill')
        srgbClr = SubElement(solidFill, 'a:srgbClr', val=border_color)
        prstDash = SubElement(ln, 'a:prstDash', val='solid')
        round_ = SubElement(ln, 'a:round')
        headEnd = SubElement(ln, 'a:headEnd', type='none', w='med', len='med')
        tailEnd = SubElement(ln, 'a:tailEnd', type='none', w='med', len='med')

def plot_PriorityvsRisk_slide1():

	slide = prs.slides.add_slide(prs.slide_layouts[3])

	#Declare the chart

	chart_data = BubbleChartData()

	#annotate the points

	series_1 = chart_data.add_series(' Category of KRI ')
	series_1.add_data_point(7.0, 3.5, 2)
	series_1.add_data_point(15.0, 3.5, 2)
	series_1.add_data_point(5.0, 2.2, 2)
	series_1.add_data_point(10.0, 2.2, 2)
	series_1.add_data_point(15.0, 2.2, 2)



	x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
	chart = slide.shapes.add_chart(
	    XL_CHART_TYPE.BUBBLE, x, y, cx, cy, chart_data
	).chart

	category_axis = chart.category_axis
	value_axis = chart.value_axis
	value_axis.axis_title.text_frame.text = 'Priority'
	category_axis.axis_title.text_frame.text = 'Risk'
	chart.has_legend = True
	value_axis.has_major_gridlines = False
	value_axis.major_tick_mark = XL_TICK_MARK.NONE
	plot  = chart.plots[0]
	# print plot.has_data_labels
	# print plot.has_data_labels
	plot.has_data_labels = True


	# assert plot.data_labels.font.fill.fore_color.type == MSO_COLOR_TYPE.RGB


	# print plot.data_labels.font.color.rgb
	tick_labels = value_axis.tick_labels
	tick_labels.number_format = '0'
	tick_labels.font.bold = True
	tick_labels.font.size = Pt(1)


	tick_labels = category_axis.tick_labels
	tick_labels.number_format = '0'
	tick_labels.font.bold = True
	tick_labels.font.size = Pt(1)
	prs.save('chart-01.pptx')
	pass

def plot_titleslide():
	# slide = prs.slides.add_slide(prs.slide_layouts[0])
	# shapes = slide.shapes
	# title = shapes.title
	# text_frame = title.text_frame
	# p = text_frame.paragraphs[0]
	# run = p.add_run()
	# run.text = 'Executive Summary – Infrastructure Risk Overview '
	# font = run.font
	# font.name = 'Verdana'
	# font.size = Pt(15)
	# font.bold = True
	# font.italic = None
	pass

def plot_executivesummary():
	TITLE_AND_CONTENT = 1
	slide = prs.slides.add_slide(prs.slide_layouts[TITLE_AND_CONTENT])
	shapes = slide.shapes
	title = shapes.title
	text_frame = title.text_frame
	p = text_frame.paragraphs[0]
	run = p.add_run()
	run.text = 'Executive Summary – Infrastructure Risk Overview '
	font = run.font
	font.name = 'Verdana'
	font.size = Pt(15)
	font.bold = True
	font.italic = None
	for shape in slide.shapes:
		if shape.is_placeholder:
			phf = shape.placeholder_format
	        print('%d, %s' % (phf.idx, phf.type))


	rows = 17
	cols = 6
	left =  Inches(0.5)
	top = Inches(0.70)
	# right = Inches(0.5)
	width = Inches(6.0)
	height = Inches(0.02)
	table = shapes.add_table(rows, cols, left, top, width, height).table

	# set column widths
	table.columns[0].width = Inches(1.5)
	table.columns[1].width = Inches(1.5)
	table.columns[2].width = Inches(3.0)

	table.cell(0, 0).text = 'ID'
	table.cell(0, 1).text = 'Category'
	table.cell(0, 2).text = 'Highlights'
	table.cell(0, 3).text = ' '
	table.cell(0, 4).text = ' '
	table.cell(0, 5).text = 'Trend'

	format_table_headers(table)
	fixed_table_content_es(table)
	format_table_content_es(table)
	prs.save('chart-01.pptx')




def plot_DetailedMetrics():
	TITLE_AND_CONTENT = 1
	slide = prs.slides.add_slide(prs.slide_layouts[TITLE_AND_CONTENT])
	shapes = slide.shapes
	title = shapes.title
	# title.text = "Detailed Metrics - KRIs "
	text_frame = title.text_frame
	# print text_frame.font.size
	p = text_frame.paragraphs[0]
	run = p.add_run()
	run.text = 'Detailed Metrics - KRIs '
	font = run.font
	font.name = 'Verdana'
	font.size = Pt(20)
	font.bold = True
	font.italic = None  # cause value to be inherited from theme



	#adding a table


	rows = 14
	cols = 6
	left =  Inches(0.5)
	top = Inches(1.10)
	# right = Inches(0.5)
	width = Inches(6.0)
	height = Inches(0.03)

	table = shapes.add_table(rows, cols, left, top, width, height).table

	# set column widths
	table.columns[0].width = Inches(1.5)
	table.columns[1].width = Inches(1.5)
	table.columns[2].width = Inches(3.0)





	# write column headings
	table.cell(0, 0).text = 'Category of KRI '
	table.cell(0, 1).text = 'Defintion'
	table.cell(0, 2).text = 'Measure'
	table.cell(0, 3).text = 'Previous Score'
	table.cell(0, 4).text = 'Current Score'
	table.cell(0, 5).text = 'Risk Trend'



	format_table_headers(table)
	cols = 6
	row = 1 
	for r in range(rows):
		for c in range(cols):
			_set_cell_border(table.cell(r,c))
			table.cell(r,c).fill.solid()
			table.cell(r,c).fill.fore_color.rgb = RGBColor(0xff, 0xff, 0xff)

	fixed_table_content(table)
	format_table_content(table)
	format_table_specific_content(table)

	for cell in iter_cells(table):
		for paragraph in cell.text_frame.paragraphs:
			for run in paragraph.runs:
				run.font.size = Pt(10)
				run.font.color.rgb = RGBColor(0,0,0)



	# _set_cell_border(table.cell(0,0))
	# _set_cell_border(table.cell(1,0))
	# _set_cell_border(table.cell(2,0))

	rule_based_formatting(table)
	row = 1 
	for row in range(1,13):
		color_in_grey(table,row,5,5)
	prs.save('chart-01.pptx')
	pass

def plot_DetailedMetrics_2():
	TITLE_AND_CONTENT = 1
	slide = prs.slides.add_slide(prs.slide_layouts[TITLE_AND_CONTENT])
	shapes = slide.shapes
	title = shapes.title
	# title.text = "Detailed Metrics - KRIs "
	text_frame = title.text_frame
	# print text_frame.font.size
	p = text_frame.paragraphs[0]
	run = p.add_run()
	run.text = 'Detailed Metrics - KRIs '
	font = run.font
	font.name = 'Verdana'
	font.size = Pt(20)
	font.bold = True
	font.italic = None  # cause value to be inherited from theme

	#adding a table


	rows = 14
	cols = 6
	left =  Inches(0.5)
	top = Inches(1.10)
	# right = Inches(0.5)
	width = Inches(6.0)
	height = Inches(0.03)

	table = shapes.add_table(rows, cols, left, top, width, height).table

	# set column widths
	table.columns[0].width = Inches(1.5)
	table.columns[1].width = Inches(1.5)
	table.columns[2].width = Inches(3.0)
	# write column headings
	table.cell(0, 0).text = 'Category of KRI '
	table.cell(0, 1).text = 'Defintion'
	table.cell(0, 2).text = 'Measure'
	table.cell(0, 3).text = 'Previous Score'
	table.cell(0, 4).text = 'Current Score'
	table.cell(0, 5).text = 'Risk Trend'

	# table.cell(1,3).text = "what is this"

	format_table_headers(table)
	cols = 6
	row = 1 
	for r in range(rows):
		for c in range(cols):
			_set_cell_border(table.cell(r,c))
			table.cell(r,c).fill.solid()
			table.cell(r,c).fill.fore_color.rgb = RGBColor(0xff, 0xff, 0xff)

	fixed_table_content_2(table)
	format_table_content(table)
	format_table_specific_content_2(table)

	for cell in iter_cells(table):
		for paragraph in cell.text_frame.paragraphs:
			for run in paragraph.runs:
				run.font.size = Pt(10)
				run.font.color.rgb = RGBColor(0,0,0)

	rule_based_formatting_2(table)

	# _set_cell_border(table.cell(0,0))
	# _set_cell_border(table.cell(1,0))
	# _set_cell_border(table.cell(2,0))
	row = 1 
	for row in range(1,14):
		color_in_grey(table,row,5,5)
	prs.save('chart-01.pptx')

	pass

def plot_DetailedMetrics_3():
	TITLE_AND_CONTENT = 1
	slide = prs.slides.add_slide(prs.slide_layouts[TITLE_AND_CONTENT])
	shapes = slide.shapes
	title = shapes.title
	# title.text = "Detailed Metrics - KRIs "
	text_frame = title.text_frame
	# print text_frame.font.size
	p = text_frame.paragraphs[0]
	run = p.add_run()
	run.text = 'Detailed Metrics - KRIs '
	font = run.font
	font.name = 'Verdana'
	font.size = Pt(20)
	font.bold = True
	font.italic = None  # cause value to be inherited from theme

	#adding a table


	rows = 15
	cols = 6
	left =  Inches(0.5)
	top = Inches(1.10)
	# right = Inches(0.5)
	width = Inches(6.0)
	height = Inches(0.03)

	table = shapes.add_table(rows, cols, left, top, width, height).table

	# set column widths
	table.columns[0].width = Inches(1.5)
	table.columns[1].width = Inches(1.5)
	table.columns[2].width = Inches(3.0)





	# write column headings
	table.cell(0, 0).text = 'Category of KRI '
	table.cell(0, 1).text = 'Defintion'
	table.cell(0, 2).text = 'Measure'
	table.cell(0, 3).text = 'Previous Score'
	table.cell(0, 4).text = 'Current Score'
	table.cell(0, 5).text = 'Risk Trend'

	# table.cell(1,3).text = "what is this"

	format_table_headers(table)
	cols = 6
	row = 1 
	for r in range(rows):
		for c in range(cols):
			_set_cell_border(table.cell(r,c))
			table.cell(r,c).fill.solid()
			table.cell(r,c).fill.fore_color.rgb = RGBColor(0xff, 0xff, 0xff)

	fixed_table_content_3(table)
	format_table_content(table)
	format_table_specific_content_3(table)

	for cell in iter_cells(table):
		for paragraph in cell.text_frame.paragraphs:
			for run in paragraph.runs:
				run.font.size = Pt(10)
				run.font.color.rgb = RGBColor(0,0,0)

	rule_based_formatting_3(table)

	# _set_cell_border(table.cell(0,0))
	# _set_cell_border(table.cell(1,0))
	# _set_cell_border(table.cell(2,0))

	row = 1 
	for row in range(1,15):
		color_in_grey(table,row,5,5)
	prs.save('chart-01.pptx')
	pass

def plot_DetailedMetrics_4():
	TITLE_AND_CONTENT = 1
	slide = prs.slides.add_slide(prs.slide_layouts[TITLE_AND_CONTENT])
	shapes = slide.shapes
	title = shapes.title
	# title.text = "Detailed Metrics - KRIs "
	text_frame = title.text_frame
	# print text_frame.font.size
	p = text_frame.paragraphs[0]

	run = p.add_run()
	run.alignment = PP_ALIGN.CENTER
	run.text = 'Detailed Metrics - KRIs '
	font = run.font
	font.name = 'Verdana'
	font.size = Pt(20)
	font.bold = True
	font.italic = None  # cause value to be inherited from theme

	#adding a table


	rows = 5
	cols = 6
	left =  Inches(0.5)
	top = Inches(1.10)
	# right = Inches(0.5)
	width = Inches(6.0)
	height = Inches(0.03)

	table = shapes.add_table(rows, cols, left, top, width, height).table

	# set column widths
	table.columns[0].width = Inches(1.5)
	table.columns[1].width = Inches(1.5)
	table.columns[2].width = Inches(3.0)





	# write column headings
	table.cell(0, 0).text = 'Category of KRI '
	table.cell(0, 1).text = 'Defintion'
	table.cell(0, 2).text = 'Measure'
	table.cell(0, 3).text = 'Previous Score'
	table.cell(0, 4).text = 'Current Score'
	table.cell(0, 5).text = 'Risk Trend'
	table.cell(0,5).alignment = PP_ALIGN.CENTER
	# table.cell(1,3).text = "what is this"

	format_table_headers(table)
	cols = 6
	row = 1 
	for r in range(rows):
		for c in range(cols):
			_set_cell_border(table.cell(r,c))
			table.cell(r,c).fill.solid()
			table.cell(r,c).fill.fore_color.rgb = RGBColor(0xff, 0xff, 0xff)

	fixed_table_content_4(table)
	format_table_content(table)
	format_table_specific_content_4(table)

	for cell in iter_cells(table):
		for paragraph in cell.text_frame.paragraphs:
			for run in paragraph.runs:
				run.font.size = Pt(10)
				run.font.color.rgb = RGBColor(0,0,0)
				run.alignment = PP_ALIGN.RIGHT

	rule_based_formatting_4(table)

	# _set_cell_border(table.cell(0,0))
	# _set_cell_border(table.cell(1,0))
	# _set_cell_border(table.cell(2,0))

	row = 1 
	for row in range(1,5):
		color_in_grey(table,row,5,5)
	prs.save('chart-01.pptx')
	pass

def plot_appendix():
	prs = pptx.Presentation('chart-01.pptx')

	prs.slide_width = 9144000


	pic_left  = int(prs.slide_width * 0.001)
	pic_top   = int(prs.slide_height * 0.01)
	pic_width = int(prs.slide_width)

	for g in glob.glob("python/Capture.PNG"):
	    print g
	    slide = prs.slides.add_slide(prs.slide_layouts[6])

	    tb = slide.shapes.add_textbox(0, 0, prs.slide_width, pic_top / 2)
	    # p = tb.textframe.add_paragraph()
	    # p.text = g
	    # p.font.size = pptx.util.Pt(14)

	    img = scipy.misc.imread(g)
	    pic_height = int(pic_width * img.shape[0] / img.shape[1])
	    #pic   = slide.shapes.add_picture(g, pic_left, pic_top)
	    pic   = slide.shapes.add_picture(g, pic_left, pic_top, pic_width, pic_height)
	prs.save('chart-01.pptx')
	pass

def plot_context():
	prs = pptx.Presentation('chart-01.pptx')

	prs.slide_width = 9144000


	pic_left  = int(prs.slide_width * 0.001)
	pic_top   = int(prs.slide_height * 0.01)
	pic_width = int(prs.slide_width)

	for g in glob.glob("python/Capture2.PNG"):
	    print g
	    slide = prs.slides.add_slide(prs.slide_layouts[6])

	    tb = slide.shapes.add_textbox(0, 0, prs.slide_width, pic_top / 2)
	    # p = tb.textframe.add_paragraph()
	    # p.text = g
	    # p.font.size = pptx.util.Pt(14)

	    img = scipy.misc.imread(g)
	    pic_height = int(pic_width * img.shape[0] / img.shape[1])
	    #pic   = slide.shapes.add_picture(g, pic_left, pic_top)
	    pic   = slide.shapes.add_picture(g, pic_left, pic_top, pic_width, pic_height)
	prs.save('chart-01.pptx')
	pass

def plot_metricdef():
	prs = pptx.Presentation('chart-01.pptx')

	prs.slide_width = 9144000


	pic_left  = int(prs.slide_width * 0.001)
	pic_top   = int(prs.slide_height * 0.01)
	pic_width = int(prs.slide_width)

	for g in glob.glob("python/Capture3.PNG"):
	    print g
	    slide = prs.slides.add_slide(prs.slide_layouts[6])

	    tb = slide.shapes.add_textbox(0, 0, prs.slide_width, pic_top / 2)
	    # p = tb.textframe.add_paragraph()
	    # p.text = g
	    # p.font.size = pptx.util.Pt(14)

	    img = scipy.misc.imread(g)
	    pic_height = int(pic_width * img.shape[0] / img.shape[1])
	    #pic   = slide.shapes.add_picture(g, pic_left, pic_top)
	    pic   = slide.shapes.add_picture(g, pic_left, pic_top, pic_width, pic_height)
	prs.save('chart-01.pptx')
	pass

def plot_threshold():
	prs = pptx.Presentation('chart-01.pptx')

	prs.slide_width = 9144000


	pic_left  = int(prs.slide_width * 0.001)
	pic_top   = int(prs.slide_height * 0.01)
	pic_width = int(prs.slide_width)

	for g in glob.glob("python/Capture4.PNG"):
	    print g
	    slide = prs.slides.add_slide(prs.slide_layouts[6])

	    tb = slide.shapes.add_textbox(0, 0, prs.slide_width, pic_top / 2)
	    # p = tb.textframe.add_paragraph()
	    # p.text = g
	    # p.font.size = pptx.util.Pt(14)

	    img = scipy.misc.imread(g)
	    pic_height = int(pic_width * img.shape[0] / img.shape[1])
	    #pic   = slide.shapes.add_picture(g, pic_left, pic_top)
	    pic   = slide.shapes.add_picture(g, pic_left, pic_top, pic_width, pic_height)
	prs.save('chart-01.pptx')
	pass

def plot_comments():
	prs = pptx.Presentation('chart-01.pptx')

	prs.slide_width = 9144000


	pic_left  = int(prs.slide_width * 0.001)
	pic_top   = int(prs.slide_height * 0.01)
	pic_width = int(prs.slide_width)

	for g in glob.glob("python/Capture5.PNG"):
	    print g
	    slide = prs.slides.add_slide(prs.slide_layouts[6])

	    tb = slide.shapes.add_textbox(0, 0, prs.slide_width, pic_top / 2)
	    # p = tb.textframe.add_paragraph()
	    # p.text = g
	    # p.font.size = pptx.util.Pt(14)

	    img = scipy.misc.imread(g)
	    pic_height = int(pic_width * img.shape[0] / img.shape[1])
	    #pic   = slide.shapes.add_picture(g, pic_left, pic_top)
	    pic   = slide.shapes.add_picture(g, pic_left, pic_top, pic_width, pic_height)
	prs.save('chart-01.pptx')
	pass



def SubElement(parent, tagname, **kwargs):
        element = OxmlElement(tagname)
        element.attrib.update(kwargs)
        parent.append(element)
        return element

def _set_cell_border(cell, border_color="000000", border_width='12700'):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for lines in ['a:lnL','a:lnR','a:lnT','a:lnB']:
        ln = SubElement(tcPr, lines, w=border_width, cap='flat', cmpd='sng', algn='ctr')
        solidFill = SubElement(ln, 'a:solidFill')
        srgbClr = SubElement(solidFill, 'a:srgbClr', val=border_color)
        prstDash = SubElement(ln, 'a:prstDash', val='solid')
        round_ = SubElement(ln, 'a:round')
        headEnd = SubElement(ln, 'a:headEnd', type='none', w='med', len='med')
        tailEnd = SubElement(ln, 'a:tailEnd', type='none', w='med', len='med')

def main():
	#create slide 1
	connect()
	populate_data()
	print "ok"
	plot_PriorityvsRisk_slide1()
	print "ok"
	plot_titleslide()
	print "ok"
	plot_executivesummary()
	print "ok"

	plot_DetailedMetrics()
	print "ok"
	plot_DetailedMetrics_2()
	print "ok"
	plot_DetailedMetrics_3()
	print "ok"
	plot_DetailedMetrics_4()
	print "ok"
	plot_appendix()
	print "ok"
	plot_context()
	print "ok"
	plot_metricdef()
	print "ok"
	plot_threshold()
	print "ok"
	plot_comments()
	print "ok"

if __name__ == "__main__": main()