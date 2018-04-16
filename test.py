from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.chart.data import BubbleChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
from pptx.chart.axis import MajorGridlines
from pptx.util import Inches, Pt, Cm, Px, Mm
from pptx.chart.data import ChartData
from pptx.chart.data import XyChartData
from pptx.chart.data import BubbleChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.enum.chart import XL_TICK_MARK
from pptx.dml.color import RGBColor
from pptx.dml.color import ColorFormat
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.oxml.xmlchemy import OxmlElement
import enum
import mysql.connector
from mysql.connector import errorcode
import pandas as pd
from decimal import Decimal
import pptx
import pptx.util
import glob
import scipy.misc

prs = Presentation()
# slide = prs.slides.add_slide(prs.slide_layouts[1])
# for shape in slide.placeholders:
#     print('%d %s' % (shape.placeholder_format.idx, shape.name))
# print slide.placeholders[0].name


slide = prs.slides.add_slide(prs.slide_layouts[1])
table_placeholder = slide.placeholders[1]
table = table_placeholder.insert_table(rows=2, cols=2)
# placeholder = slide.placeholders[10]  # idx key, not position
# print placeholder.name
# # //'Table Placeholder 1'
# print placeholder.shape_type
# //TABLE (12)
# graphic_frame = placeholder.insert_table(rows=2, cols=2)
# table = graphic_frame.table
# print  len(table.rows), len(table.columns)
# //(2, 2)
