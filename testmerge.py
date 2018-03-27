from pptx import Presentation
from pptx.util import Inches,Pt

prs = Presentation()
title_only_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_only_slide_layout)
shapes = slide.shapes
rows = 6
cols = 6
left = top = Inches(0.5)
right = Inches(0.5)
width = Inches(5.0)
height = Inches(0.2)

table = shapes.add_table(rows, cols, left, top, width, height).table


# d  = { }
# d['Inventory']={'Score' : '45' , 'Previous Score' : '34'}
# print d['Inventory'].keys()
# for element in d.keys():
# 	table.cell(1,0).text = element
# 	for item in d[element]:
# 		table.cell(1,1).text = "45"
# 		table.cell(1,2)


dict = { }
arr = [[1,2,3,4,5],[1,2,3,0,0],[1,5,0,0,0],[0,0,1,2,3],[1,2,0,9,0]]
dict = { 'Cell Values' :  arr  }
print dict.keys()
if 'Cell Values' in dict.keys():
	for j in range(len(dict['Cell Values'])):
		for i in range(len(dict['Cell Values'][j])):
			table.cell(i, j).text = str(dict['Cell Values'][i][j])
			if dict['Cell Values'][i][j] > 0:
				row_cells = []
				row_cells.append(table.cell(i, j))
			elif dict['Cell Values'][i][j] == 0:
				print "Found empty cell"
				row_cells.append(table.cell(i, j))
				row_cells[0]._tc.set('gridSpan',str(len(row_cells)))
				row_cells[-1]._tc.set('hMerge', '1')

else:
	print "Cell Values NOT FOUND"
	table.cell(0, 0).text = 'Title Row Col 1'
	table.cell(0, 1).text = 'Title Row Col 2'
	table.cell(1, 0).text = 'Body Row Col 1'
	table.cell(1, 1).text = 'Body Row Col 2'


prs.save('testmerge.pptx')
