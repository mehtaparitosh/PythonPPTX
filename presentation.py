from pptx import Presentation


def iter_cells(table):
    for row in table.rows:
        for cell in row.cells:
            yield cell

def Create_Table_On_Slide(self, shapes, Table_Dict={}):
		if 'Cell Values' in Table_Dict.keys():
			rows = len(Table_Dict['Cell Values'])
			cols = len(Table_Dict['Cell Values'][0])
		else:
			rows = 2
			cols = 2

		if 'Left' in Table_Dict.keys():
			left = Inches(Table_Dict['Left'])
		else:
			left = Inches(2.0)

		if 'Top' in Table_Dict.keys():
			top = Inches(Table_Dict['Top'])
		else:
			top = Inches(2.0)
		
		if 'Row Height' in Table_Dict.keys():
			rowHeight = Table_Dict['Row Height']
		else:
			rowHeight = 0.4
		height = Inches(rows * rowHeight)

		if 'Col Widths' in Table_Dict.keys():
			width = Inches(sum(Table_Dict['Col Widths']))

		table = shapes.add_table(rows, cols, left, top, width, height).table

		# set column widths
		if 'Col Widths' in Table_Dict.keys():
			for i in range(len(Table_Dict['Col Widths'])):
				table.columns[i].width = Inches(Table_Dict['Col Widths'][i])
		else:
			table.columns[0].width = Inches(2.0)
			table.columns[1].width = Inches(4.0)

		# write column headings and body cells
		# To merge cells: https://groups.google.com/forum/#!topic/python-pptx/cVRP9sSpEjA
		if 'Cell Values' in Table_Dict.keys():
			for i in range(len(Table_Dict['Cell Values'])):
				for j in range(len(Table_Dict['Cell Values'][i])):
					table.cell(i, j).text = str(Table_Dict['Cell Values'][i][j])
					if len(table.cell(i, j).text_frame.text) > 0:
						row_cells = []
						row_cells.append(table.cell(i, j))
					elif len(table.cell(i, j).text_frame.text) == 0:
						row_cells.append(table.cell(i, j))
						row_cells[0]._tc.set('gridSpan',str(len(row_cells)))
						row_cells[-1]._tc.set('hMerge', '1')

		else:
			table.cell(0, 0).text = 'Title Row Col 1'
			table.cell(0, 1).text = 'Title Row Col 2'
			table.cell(1, 0).text = 'Body Row Col 1'
			table.cell(1, 1).text = 'Body Row Col 2'



def main():

	prs = Presentation()
	title_slide_layout = prs.slide_layouts[0]
	slide = prs.slides.add_slide(title_slide_layout)
	title = slide.shapes.title
	subtitle = slide.placeholders[1]

	title.text = "Hello, World!"
	subtitle.text = "python-pptx was here!"



	prs.save('test.pptx')


	# prs = Presentation()
	bullet_slide_layout = prs.slide_layouts[1]

	slide = prs.slides.add_slide(bullet_slide_layout)
	shapes = slide.shapes

	title_shape = shapes.title
	body_shape = shapes.placeholders[1]

	title_shape.text = 'Adding a Bullet Slide'

	tf = body_shape.text_frame
	tf.text = 'Find the bullet slide layout'

	p = tf.add_paragraph()
	p.text = 'Use _TextFrame.text for first bullet'
	p.level = 1

	p = tf.add_paragraph()
	p.text = 'Use _TextFrame.add_paragraph() for subsequent bullets'
	p.level = 2



	from pptx.util import Inches,Pt

	# prs = Presentation()
	title_only_slide_layout = prs.slide_layouts[5]
	slide = prs.slides.add_slide(title_only_slide_layout)
	shapes = slide.shapes

	shapes.title.text = 'Adding a Table'

	rows = 15
	cols = 6
	left = top = Inches(0.5)
	right = Inches(0.5)
	width = Inches(5.0)
	height = Inches(0.2)

	table = shapes.add_table(rows, cols, left, top, width, height).table

	# set column widths
	table.columns[0].width = Inches(2.0)
	table.columns[1].width = Inches(4.0)

	# write column headings
	table.cell(0, 0).text = 'Category of KRI '
	table.cell(0, 1).text = 'Defintion'
	table.cell(0, 2).text = 'Measure'
	table.cell(0, 3).text = 'Previous Score'
	table.cell(0, 4).text = 'Current Score'
	table.cell(0, 5).text = 'Risk Trend'



	# write body cells
	# table.cell(2, 0).text = 'Inventory'
	for pp in table.cell(1,0).text_frame.paragraphs:
		for run in pp.runs:
			run.font.bold = True 

	table.cell(1, 1).text = 'Qux'

	for cell in iter_cells(table):
		for paragraph in cell.text_frame.paragraphs:
			for run in paragraph.runs:
				run.font.size = Pt(12)

	
	table.cell(1,0).text = 'Random'
	table.cell(1,1).text = 'Random'
	table.cell(1,2).text = 'Random'
	table.cell(1,3).text = 'Random'
	table.cell(1,4).text = 'Random'

	# table.cell(2,0).text = 'Inventory'
	table.cell(5,0).text = 'random'
	table.cell(6,0).text = 'Random'
	table.cell(7,0).text = 'Random'

	col_cells = []



	col_cells.append(table.cell(1,0))
	# row_cells.append(table.cell(2,0))
	# col_cells.append(table.cell(5,0))
	# col_cells.append(table.cell(6,0))
	# col_cells.append(table.cell(7,0))

	print table.cell(2,0).text_frame.text
	print table.cell(0,2).text_frame.text


	if len(table.cell(3,0).text_frame.text) > 0:
		row_cells.append(table.cell(1,0))
	elif len(table.cell(3,0).text_frame.text) == 0:
		print "IS empty bruv"
		# row_cells.append(table.cell(2,0))
		col_cells[0]._tc.set('gridSpan',str(len(col_cells)))
		col_cells[-1]._tc.set('vMerge', '1')
		


	vbcode = """
	
	With ActivePresentation.Slides(3).Shapes(5).Table

    .Cell(3, 2).Merge MergeTo:=.Cell(4, 2)

	End With
"""	



	# if len(table.cell(2,1).text_frame.text) > 0:
	# 	row_cells.append(table.cell(1,0))
	# elif len(table.cell(2,1).text_frame.text) == 0:
	# 	print "IS empty bruv"
	# 	# row_cells.append(table.cell(2,0))
	# 	row_cells[1]._tc.set('gridSpan',str(len(row_cells)))
	# 	row_cells[0]._tc.set('vMerge', '1')	
	# prs.save('test.pptx')

	# if len(table.cell(2,2).text_frame.text) > 0:
	# 	row_cells.append(table.cell(1,0))
	# elif len(table.cell(2,2).text_frame.text) == 0:
	# 	print "IS empty bruv"
	# 	# row_cells.append(table.cell(2,0))
	# 	row_cells[2]._tc.set('gridSpan',str(len(row_cells)))
	# 	row_cells[1]._tc.set('vMerge', '1')	
	# # prs.save('test.pptx')
	prs.save('test.pptx')



		# from pptx import Presentation
	

if __name__ == "__main__": main()